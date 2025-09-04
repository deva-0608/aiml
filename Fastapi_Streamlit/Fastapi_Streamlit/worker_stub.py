import time
import json
from pathlib import Path
from uni import run_uvi   # our wrapper for data.py + insights
from pptx import Presentation


def create_ppt_and_preview(description_json, column_separation, out_dir, job_id):
    """Generate a simple PPT and preview.json from the analysis results"""
    prs = Presentation()
    preview = {"job_id": job_id, "slides": []}

    # --- Slide 1: Title ---
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "Dataset Analysis Report"
    slide.placeholders[1].text = f"Job ID: {job_id}"
    preview["slides"].append({
        "title": "Dataset Analysis Report",
        "content": f"Job ID: {job_id}"
    })

    # --- Slide 2: Dataset Overview ---
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Dataset Overview"
    summary = description_json.get("summary", {})
    content = slide.shapes.placeholders[1].text_frame
    for k, v in summary.items():
        content.add_paragraph(f"{k}: {v}")
    preview["slides"].append({
        "title": "Dataset Overview",
        "content": summary
    })

    # --- Slide 3: Columns Overview ---
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Columns Overview"
    content = slide.shapes.placeholders[1].text_frame
    content.add_paragraph(f"Numerical: {len(column_separation.get('numerical', []))}")
    content.add_paragraph(f"Categorical: {len(column_separation.get('categorical', []))}")
    content.add_paragraph(f"Datetime: {len(column_separation.get('datetime', []))}")
    preview["slides"].append({
        "title": "Columns Overview",
        "content": column_separation
    })

    # Save PPT + Preview JSON
    ppt_path = out_dir / f"{job_id}.pptx"
    prs.save(ppt_path)

    preview_path = out_dir / "preview.json"
    with open(preview_path, "w") as f:
        json.dump(preview, f, indent=2)

    return ppt_path, preview_path


   # <-- use your updated uvi.py



# Directories
UPLOADS_DIR = Path("storage/uploads")
OUTPUTS_DIR = Path("storage/outputs")


def worker_loop(poll_interval=5):
    """Worker that continuously watches for new uploads and runs UVI pipeline"""
    print("👷 Worker started. Watching for uploads...")

    # Ensure base folders exist
    UPLOADS_DIR.mkdir(parents=True, exist_ok=True)
    OUTPUTS_DIR.mkdir(parents=True, exist_ok=True)

    while True:
        try:
            for job_dir in UPLOADS_DIR.glob("*"):
                if not job_dir.is_dir():
                    continue

                job_id = job_dir.name
                out_dir = OUTPUTS_DIR / job_id
                out_dir.mkdir(parents=True, exist_ok=True)

                status_path = out_dir / "status.txt"

                # Skip if already completed
                if (out_dir / "description.json").exists() and (out_dir / "insights.json").exists() and status_path.exists():
                    continue

                # Look for supported files
                files = list(job_dir.glob("*.csv")) + list(job_dir.glob("*.xlsx")) + list(job_dir.glob("*.xls"))
                if not files:
                    continue

                input_file = str(files[0])
                print(f"📂 Found file for job {job_id}: {input_file}")

                # Mark status as processing
                with open(status_path, "w") as f:
                    f.write("processing")

                # Run UVI pipeline
                try:
                    run_uvi(input_file, job_id=job_id, output_dir=str(OUTPUTS_DIR))

                    # Update status
                    with open(status_path, "w") as f:
                        f.write("completed")

                    print(f"✅ Job {job_id} completed successfully")
                    print(f"   → description.json & insights.json generated")

                except Exception as e:
                    with open(status_path, "w") as f:
                        f.write("failed")

                    print(f"❌ Job {job_id} failed: {e}")

        except Exception as e:
            print(f"❌ Worker loop error: {str(e)}")

        time.sleep(poll_interval)


if __name__ == "__main__":
    worker_loop()