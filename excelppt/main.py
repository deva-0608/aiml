# main.py
# FastAPI backend for Excel/CSV upload, job status, preview, and download.
# Your parsing/RulePy/PPT generation should create:
#   storage/outputs/<job_id>/presentation.pptx
#   storage/outputs/<job_id>/preview.json   (optional but recommended)
#
# Run: uvicorn main:app --reload --port 8000

import os
import io
import json
import uuid
from typing import List, Optional
from fastapi import FastAPI, UploadFile, File, HTTPException
from fastapi.responses import FileResponse, JSONResponse
from fastapi.middleware.cors import CORSMiddleware

STORAGE_ROOT = os.path.join(os.path.dirname(__file__), "storage")
UPLOADS_DIR = os.path.join(STORAGE_ROOT, "uploads")
OUTPUTS_DIR = os.path.join(STORAGE_ROOT, "outputs")
os.makedirs(UPLOADS_DIR, exist_ok=True)
os.makedirs(OUTPUTS_DIR, exist_ok=True)

app = FastAPI(title="Excel→PPT Backend")

# Allow Streamlit (same machine) to talk to backend without CORS pain
app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],  # tighten this in prod
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

def job_dirs(job_id: str):
    return (
        os.path.join(UPLOADS_DIR, job_id),
        os.path.join(OUTPUTS_DIR, job_id),
    )

def find_ppt_path(job_id: str) -> Optional[str]:
    out_dir = os.path.join(OUTPUTS_DIR, job_id)
    ppt_path = os.path.join(out_dir, "presentation.pptx")
    return ppt_path if os.path.exists(ppt_path) else None

def load_preview(job_id: str):
    """Try preview.json first. If absent, try a tiny fallback using python-pptx (optional)."""
    out_dir = os.path.join(OUTPUTS_DIR, job_id)
    manifest = os.path.join(out_dir, "preview.json")
    if os.path.exists(manifest):
        with open(manifest, "r", encoding="utf-8") as f:
            data = json.load(f)
        # expected schema:
        # { "slides": [ {"title": "...", "subtitle": "..."} , ... ] }
        return data

    # Fallback: attempt minimal text extraction if python-pptx is available
    ppt_path = find_ppt_path(job_id)
    if ppt_path is None:
        return None

    try:
        from pptx import Presentation  # type: ignore
        prs = Presentation(ppt_path)
        slides = []
        for s in prs.slides:
            title = None
            subtitle = None
            # naive text pull
            txts = []
            for shp in s.shapes:
                try:
                    if shp.has_text_frame:
                        txts.append(shp.text.strip())
                except Exception:
                    pass
            if txts:
                title = txts[0]
                if len(txts) > 1:
                    subtitle = txts[1]
            slides.append({"title": title or "Slide", "subtitle": subtitle or ""})
        return {"slides": slides}
    except Exception:
        # No python-pptx or failure: minimal placeholder
        return {"slides": [{"title": "Data Presentation", "subtitle": f"Generated for job {job_id}"}]}

@app.get("/")
async def home():
    return {"message":"the app is running succesfully"}


@app.post("/upload")
async def upload_file(file: UploadFile = File(...)):
    ext = os.path.splitext(file.filename or "")[1].lower()
    if ext not in [".csv", ".xlsx", ".xls"]:
        raise HTTPException(status_code=400, detail="Please upload a CSV or Excel file.")
    job_id = str(uuid.uuid4())[:8]
    up_dir, out_dir = job_dirs(job_id)
    os.makedirs(up_dir, exist_ok=True)
    os.makedirs(out_dir, exist_ok=True)

    # Save original file
    dest = os.path.join(up_dir, f"input{ext}")
    content = await file.read()
    with open(dest, "wb") as f:
        f.write(content)

    # ⚠️ Your pipeline should watch `storage/uploads/<job_id>` and
    # create `storage/outputs/<job_id>/presentation.pptx` (+ optional preview.json)
    # This backend does NOT do conversion; it only exposes endpoints.
    return {"job_id": job_id, "filename": file.filename}

@app.get("/status/{job_id}")
def job_status(job_id: str):
    ppt = find_ppt_path(job_id)
    if not ppt:
        return {"status": "pending", "job_id": job_id}
    # If PPT exists, try preview
    preview = load_preview(job_id) or {"slides": [{"title": "Data Presentation", "subtitle": ""}]}
    slide_count = len(preview.get("slides", []))
    return {
        "status": "ready",
        "job_id": job_id,
        "slide_count": slide_count,
        "preview": preview,
        "download_url": f"/download/{job_id}"
    }

@app.get("/preview/{job_id}")
def preview(job_id: str):
    ppt = find_ppt_path(job_id)
    if not ppt:
        return {"status": "pending"}
    return {"status": "ready", "preview": load_preview(job_id)}

@app.get("/download/{job_id}")
def download(job_id: str):
    ppt = find_ppt_path(job_id)
    if not ppt:
        raise HTTPException(status_code=404, detail="PPT not found yet.")
    filename = f"presentation_{job_id}.pptx"
    return FileResponse(ppt, media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation", filename=filename)
