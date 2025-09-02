# worker_stub.py
# Simulates a generator: takes any uploaded Excel/CSV file,
# and creates a simple PowerPoint + preview.json in the outputs folder.

import os
import time
import glob
import json
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt

BASE_DIR = Path(__file__).parent
UPLOADS_DIR = BASE_DIR / "storage" / "uploads"
OUTPUTS_DIR = BASE_DIR / "storage" / "outputs"

def make_sample_ppt(job_id: str, input_file: str):
    """Create a dummy PPT for testing with python-pptx."""
    prs = Presentation()
    # Title slide
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "Data Presentation"
    subtitle.text = f"Generated from {os.path.basename(input_file)}"

    # Add 3 simple slides
    bullet_slide_layout = prs.slide_layouts[1]
    for i in range(1, 4):
        slide = prs.slides.add_slide(bullet_slide_layout)
        shapes = slide.shapes
        shapes.title.text = f"Slide {i}"
        body_shape = shapes.placeholders[1]
        tf = body_shape.text_frame
        tf.text = f"This is auto-generated content for slide {i}."

    # Save PPT
    out_dir = OUTPUTS_DIR / job_id
    out_dir.mkdir(parents=True, exist_ok=True)
    ppt_path = out_dir / "presentation.pptx"
    prs.save(ppt_path)

    # Create preview.json
    preview = {
        "slides": [
            {"title": "Data Presentation", "subtitle": f"Generated from {os.path.basename(input_file)}"},
            {"title": "Slide 1", "subtitle": "This is auto-generated content for slide 1."},
            {"title": "Slide 2", "subtitle": "This is auto-generated content for slide 2."},
            {"title": "Slide 3", "subtitle": "This is auto-generated content for slide 3."}
        ]
    }
    with open(out_dir / "preview.json", "w", encoding="utf-8") as f:
        json.dump(preview, f, indent=2)

    print(f"✅ Generated PPT + preview.json for job {job_id}")

def worker_loop(poll_interval=5):
    print("Worker started. Watching for uploads…")
    while True:
        # Check for uploaded jobs
        for job_dir in UPLOADS_DIR.glob("*"):
            job_id = job_dir.name
            out_dir = OUTPUTS_DIR / job_id
            ppt_path = out_dir / "presentation.pptx"
            if ppt_path.exists():
                continue  # already processed

            # Pick first file in job upload dir
            files = list(job_dir.glob("*"))
            if not files:
                continue
            input_file = str(files[0])
            make_sample_ppt(job_id, input_file)

        time.sleep(poll_interval)

if __name__ == "__main__":
    worker_loop()
